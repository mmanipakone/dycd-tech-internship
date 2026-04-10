"""
PDF to PPTX Converter — TechInternship_final.pdf
Strategy:
  1. Render each PDF page at 3x resolution → PNG background (preserves all design)
  2. Extract all text spans with position, font, size, color, bold/italic
  3. Overlay transparent editable text boxes at exact positions
  4. Group spans into logical paragraphs per PDF block
Result: fully editable text over pixel-perfect page designs
"""

import fitz  # PyMuPDF
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paths ────────────────────────────────────────────────────────────────────
PDF_PATH = '/Users/manipakone/Documents/dycd-tech-internship/playbook/TechInternship_final.pdf'
OUT_PATH = '/Users/manipakone/Documents/dycd-tech-internship/TechInternship_final.pptx'

# ── Slide dimensions (8.5 × 11 in = US Letter) ───────────────────────────────
SLIDE_W = Inches(8.5)
SLIDE_H = Inches(11.0)

# PDF coordinate system: 1 pt = 1/72 inch = 12700 EMU
PT_TO_EMU = 12700

# Render zoom (higher = crisper background image)
ZOOM = 3.0

# ── Helpers ──────────────────────────────────────────────────────────────────

def pt(val):
    """Convert PDF points to EMU."""
    return int(val * PT_TO_EMU)

def color_rgb(color_int):
    """Convert PDF color int 0xRRGGBB to pptx RGBColor."""
    r = (color_int >> 16) & 0xFF
    g = (color_int >> 8) & 0xFF
    b = color_int & 0xFF
    return RGBColor(r, g, b)

def map_font(font_name: str) -> str:
    """Map PDF font names to available system fonts."""
    if not font_name:
        return 'Arial'
    n = font_name.lower()
    if 'arial' in n:
        if 'bold' in n:
            return 'Arial'
        if 'italic' in n:
            return 'Arial'
        return 'Arial'
    if 'helvetica' in n:
        return 'Arial'
    if 'symbol' in n or 'unnamed' in n:
        return 'Arial'
    return 'Arial'

def is_bold(flags: int) -> bool:
    return bool(flags & 16)

def is_italic(flags: int) -> bool:
    return bool(flags & 2)

def clamp_emu(val, max_val):
    return max(0, min(val, max_val))

def set_no_fill(shape):
    """Set shape fill to transparent (no fill)."""
    fill = shape.fill
    fill.background()

def set_no_line(shape):
    """Remove shape border."""
    shape.line.fill.background()

def move_to_back(slide, element):
    """Move a shape element behind all others (z-order = back)."""
    sp_tree = slide.shapes._spTree
    sp_tree.remove(element)
    sp_tree.insert(2, element)  # index 2 = behind all other shapes

def set_text_no_margins(tf):
    """Zero out text frame internal margins."""
    tf.margin_left  = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top   = Pt(0)
    tf.margin_bottom = Pt(0)

# ── Main conversion ───────────────────────────────────────────────────────────

def convert():
    doc = fitz.open(PDF_PATH)
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # completely blank

    slides_with_image_only = []  # pages where text extraction was skipped

    for page_num in range(len(doc)):
        page = doc[page_num]
        slide = prs.slides.add_slide(blank_layout)

        page_w = page.rect.width   # points
        page_h = page.rect.height  # points

        # ── Step 1: Render page to high-res PNG ──────────────────────────────
        mat = fitz.Matrix(ZOOM, ZOOM)
        pix = page.get_pixmap(matrix=mat, alpha=False, colorspace=fitz.csRGB)
        img_bytes = pix.tobytes("png")
        img_stream = io.BytesIO(img_bytes)

        # Add as full-slide background image
        pic = slide.shapes.add_picture(img_stream, 0, 0, SLIDE_W, SLIDE_H)
        move_to_back(slide, pic._element)

        # ── Step 2: Extract text blocks ──────────────────────────────────────
        blocks = page.get_text(
            'dict',
            flags=(fitz.TEXT_PRESERVE_WHITESPACE | fitz.TEXT_PRESERVE_SPANS)
        )['blocks']

        text_blocks = [b for b in blocks if b['type'] == 0]

        if not text_blocks:
            slides_with_image_only.append(page_num + 1)
            continue

        for block in text_blocks:
            bx0, by0, bx1, by1 = block['bbox']

            # Build flattened list of (line_idx, spans) preserving line breaks
            lines = block.get('lines', [])
            if not lines:
                continue

            # Compute text box bounds (add small padding)
            PAD = 2  # points
            left_emu  = clamp_emu(pt(bx0 - PAD), SLIDE_W)
            top_emu   = clamp_emu(pt(by0 - PAD), SLIDE_H)
            box_w_emu = clamp_emu(pt(bx1 - bx0 + PAD * 2), SLIDE_W - left_emu)
            box_h_emu = clamp_emu(pt(by1 - by0 + PAD * 2), SLIDE_H - top_emu)

            # Skip degenerate boxes
            if box_w_emu < pt(2) or box_h_emu < pt(2):
                continue

            txBox = slide.shapes.add_textbox(
                left_emu, top_emu, box_w_emu, box_h_emu
            )
            set_no_fill(txBox)
            set_no_line(txBox)

            tf = txBox.text_frame
            tf.word_wrap = False
            set_text_no_margins(tf)

            first_para = True
            for line_idx, line in enumerate(lines):
                spans = line.get('spans', [])
                if not spans:
                    continue

                if first_para:
                    para = tf.paragraphs[0]
                    first_para = False
                else:
                    para = tf.add_paragraph()

                # Para spacing: tighten up to match PDF layout
                pPr = para._pPr
                if pPr is None:
                    pPr = para._p.get_or_add_pPr()
                # Set line spacing to single (100%)
                lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
                spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
                spcPct.set('val', '100000')  # 100%
                # No space before/after
                spcBef = etree.SubElement(pPr, qn('a:spcBef'))
                spcBef_pts = etree.SubElement(spcBef, qn('a:spcPts'))
                spcBef_pts.set('val', '0')
                spcAft = etree.SubElement(pPr, qn('a:spcAft'))
                spcAft_pts = etree.SubElement(spcAft, qn('a:spcPts'))
                spcAft_pts.set('val', '0')

                for span in spans:
                    raw_text = span.get('text', '')
                    if not raw_text:
                        continue

                    run = para.add_run()
                    run.text = raw_text

                    font = run.font
                    font.size  = Pt(round(span['size'], 1))
                    font.name  = map_font(span.get('font', ''))
                    font.bold  = is_bold(span['flags'])
                    font.italic = is_italic(span['flags'])

                    color_int = span.get('color', 0)
                    font.color.rgb = color_rgb(color_int)

        print(f"  Page {page_num + 1}/{len(doc)} done")

    doc.close()
    prs.save(OUT_PATH)

    print(f"\nSaved: {OUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")
    if slides_with_image_only:
        print(f"Image-only pages (no text extracted): {slides_with_image_only}")
    else:
        print("All pages have editable text overlays.")

if __name__ == '__main__':
    convert()
