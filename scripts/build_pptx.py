"""
Build TechInternship_final.pptx — fully editable PowerPoint (no image backgrounds)
16 slides: all text, shapes, and tables as native PPT elements
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUT = '/Users/manipakone/Documents/dycd-tech-internship/TechInternship_final.pptx'

# ── Colors ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x25, 0x40)
ORANGE = RGBColor(0xE8, 0x4B, 0x00)
YELLOW = RGBColor(0xFF, 0xE6, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LGRAY  = RGBColor(0xE1, 0xE8, 0xEF)
STRIPE = RGBColor(0xF2, 0xF5, 0xF9)
LNAVY  = RGBColor(0x14, 0x3A, 0x63)

# ── Dimensions ────────────────────────────────────────────────────────────────
W  = Inches(8.5)
H  = Inches(11.0)
ML = Inches(0.5)
MR = Inches(0.5)
CW = W - ML - MR       # 7.5"

HDR_H  = Inches(0.35)
FTR_Y  = Inches(10.65)
FTR_H  = Inches(0.22)
BODY_Y = HDR_H + Inches(0.12)

C1W = Inches(3.45)
GAP = Inches(0.18)
C2X = ML + C1W + GAP
C2W = W - C2X - MR    # ~3.82"

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank

def rect(s, x, y, w, h, fill, line_color=None):
    shape = s.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def txt(s, x, y, w, h, text, size=10, color=DARK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, wrap=True, va='top'):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    if va == 'middle':
        tf.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb

def mtxt(s, x, y, w, h, paras, wrap=True):
    """Multi-paragraph text box. paras = list of dicts with keys:
       text, size, color, bold, italic, align, space_before, bullet, indent
    """
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    first = True
    for p_def in paras:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = p_def.get('align', PP_ALIGN.LEFT)
        # space before
        sb = p_def.get('space_before', 0)
        if sb:
            pPr = p._p.get_or_add_pPr()
            spcBef = etree.SubElement(pPr, qn('a:spcBef'))
            spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
            spcPts.set('val', str(int(sb * 100)))
        # indent / bullet
        indent = p_def.get('indent', 0)
        if indent:
            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', str(int(Inches(indent * 0.18))))
            pPr.set('indent', str(int(-Inches(0.18))))
        for seg in p_def.get('runs', [p_def]):
            r = p.add_run()
            r.text = seg.get('text', '')
            r.font.size = Pt(seg.get('size', 10))
            r.font.color.rgb = seg.get('color', DARK)
            r.font.bold = seg.get('bold', False)
            r.font.italic = seg.get('italic', False)
    return tb

def hdr_ftr(s, section, page_num, footer_label=None):
    """Navy header bar + footer bar."""
    # Header
    r = rect(s, 0, 0, W, HDR_H, NAVY)
    # Section label in header
    txt(s, ML, Pt(4), Inches(6), HDR_H,
        section, size=8, color=WHITE, bold=True)
    # Page number
    txt(s, W - Inches(1), Pt(4), Inches(0.9), HDR_H,
        str(page_num).zfill(2), size=9, color=ORANGE, bold=True, align=PP_ALIGN.RIGHT)
    # Footer
    rect(s, 0, FTR_Y, W, FTR_H, NAVY)
    foot = footer_label or f'NYC DYCD — High School Tech Internship Framework'
    txt(s, ML, FTR_Y + Pt(2), W - ML - MR, FTR_H,
        foot, size=7, color=WHITE, align=PP_ALIGN.CENTER)

def label_bar(s, x, y, w, text, fill=NAVY, text_color=WHITE, size=8.5, h=Inches(0.24)):
    """Colored label/section bar."""
    rect(s, x, y, w, h, fill)
    txt(s, x + Pt(6), y + Pt(2), w - Pt(8), h,
        text, size=size, color=text_color, bold=True)
    return y + h

def cell_fill(cell, color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(old)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    clr = etree.SubElement(sf, qn('a:srgbClr'))
    clr.set('val', str(color))  # RGBColor is a str subclass, str() gives 6-char hex

def cell_fmt(cell, text, size=9, color=DARK, bold=False, fill=None,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    if fill:
        cell_fill(cell, fill)
    tf = cell.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    # clear existing runs
    for r in p.runs:
        r.text = ''
    if p.runs:
        r = p.runs[0]
    else:
        r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic

def set_col_widths(table, widths):
    """widths: list of Inches/Emu values."""
    tbl_elem = table._tbl
    tblGrid = tbl_elem.find(qn('a:tblGrid'))
    cols = tblGrid.findall(qn('a:gridCol'))
    for i, col in enumerate(cols):
        if i < len(widths):
            col.set('w', str(int(widths[i])))

def make_table(s, x, y, w, h, headers, rows, col_widths=None,
               hdr_fill=NAVY, hdr_color=WHITE, hdr_size=8.5,
               row_size=8.5, stripe=True):
    """Create a formatted table."""
    nrows = len(rows) + 1
    ncols = len(headers)
    tbl = s.shapes.add_table(nrows, ncols, x, y, w, h).table
    tbl.first_row = True

    if col_widths:
        set_col_widths(tbl, col_widths)

    # Header row
    for j, hdr in enumerate(headers):
        cell_fmt(tbl.cell(0, j), hdr, size=hdr_size,
                 color=hdr_color, bold=True, fill=hdr_fill,
                 align=PP_ALIGN.LEFT)

    # Data rows
    for i, row in enumerate(rows):
        bg = STRIPE if (stripe and i % 2 == 1) else WHITE
        for j, val in enumerate(row):
            if isinstance(val, dict):
                cell_fmt(tbl.cell(i+1, j),
                         val.get('text', ''),
                         size=val.get('size', row_size),
                         color=val.get('color', DARK),
                         bold=val.get('bold', False),
                         fill=val.get('fill', bg),
                         align=val.get('align', PP_ALIGN.LEFT),
                         italic=val.get('italic', False))
            else:
                cell_fmt(tbl.cell(i+1, j), str(val),
                         size=row_size, fill=bg)
    return tbl

# ── Slide 1: Cover ────────────────────────────────────────────────────────────

def slide1(prs):
    s = new_slide(prs)
    # Full navy background
    rect(s, 0, 0, W, H, NAVY)
    # Orange accent bar at top
    rect(s, 0, 0, W, Inches(0.06), ORANGE)
    # Orange accent bar at bottom
    rect(s, 0, H - Inches(0.06), W, Inches(0.06), ORANGE)

    # "TECH INTERNSHIP" — large orange
    txt(s, ML, Inches(1.8), CW, Inches(0.8),
        'TECH INTERNSHIP', size=42, color=ORANGE, bold=True)
    # "Summer Youth Employment Program (SYEP)"
    txt(s, ML, Inches(2.65), CW, Inches(0.4),
        'Summer Youth Employment Program (SYEP)', size=13, color=WHITE)
    # Divider line
    rect(s, ML, Inches(3.12), CW, Pt(1.5), ORANGE)
    # "PROVIDER PLAYBOOK"
    txt(s, ML, Inches(3.2), CW, Inches(0.7),
        'PROVIDER PLAYBOOK', size=30, color=WHITE, bold=True)
    # "6-WEEK FRAMEWORK"
    txt(s, ML, Inches(3.95), CW, Inches(0.35),
        '6-WEEK FRAMEWORK', size=12, color=ORANGE, bold=True)

    # Badge: "3 PROJECT SCENARIOS INCLUDED"
    rect(s, ML, Inches(4.45), Inches(3.2), Inches(0.3), ORANGE)
    txt(s, ML + Pt(6), Inches(4.47), Inches(3.0), Inches(0.28),
        '3 PROJECT SCENARIOS INCLUDED', size=9, color=WHITE, bold=True)

    # "PLAYBOOK" label
    rect(s, ML, Inches(5.0), Inches(1.5), Inches(0.26), LNAVY)
    txt(s, ML + Pt(6), Inches(5.02), Inches(1.3), Inches(0.24),
        'PLAYBOOK', size=8, color=WHITE, bold=True)

    # Description
    txt(s, ML, Inches(5.35), CW, Inches(0.7),
        'A comprehensive implementation guide for employers hosting high school interns\n'
        'in technology-focused internship experiences.',
        size=10, color=LGRAY)

    # "HIGH SCHOOL" badge
    rect(s, ML, Inches(6.2), Inches(1.3), Inches(0.26), ORANGE)
    txt(s, ML + Pt(6), Inches(6.22), Inches(1.2), Inches(0.24),
        'HIGH SCHOOL', size=8, color=WHITE, bold=True)

    # DYCD text
    txt(s, ML, Inches(9.8), CW, Inches(0.3),
        'Department of Youth and Community Development', size=8, color=LGRAY)
    txt(s, ML, Inches(10.1), CW, Inches(0.3),
        'New York City Department of Youth & Community Development',
        size=8, color=GRAY)

    print('  Slide 1 done')

# ── Slide 2: Table of Contents ────────────────────────────────────────────────

def slide2(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'Table of contents', 2)
    rect(s, 0, 0, W, HDR_H, NAVY)

    y = BODY_Y

    # "What's inside" heading
    txt(s, ML, y, CW, Inches(0.4),
        "What's inside", size=18, color=NAVY, bold=True)
    y += Inches(0.42)

    # Description
    txt(s, ML, y, CW, Inches(0.5),
        'This playbook equips host employers and supervisors with everything needed to run a '
        'high-quality, skills-driven tech internship experience for NYC high school students through SYEP.',
        size=8.5, color=GRAY, wrap=True)
    y += Inches(0.55)

    # Goal callout
    rect(s, ML, y, CW, Inches(0.35), STRIPE)
    txt(s, ML + Pt(8), y + Pt(4), CW - Pt(10), Inches(0.3),
        '\u26a1 Program Goal: Build transferable tech skills while giving NYC youth real-world '
        'professional experience in a supportive, structured environment.',
        size=8.5, color=NAVY, bold=False)
    y += Inches(0.42)

    # TOC items — two columns
    left_items = [
        ('3',  'Why Host + Portrait of an Intern'),
        ('4',  'The 6-Week Internship Journey'),
        ('5',  'Weekly Schedule & Time Breakdown'),
        ('6',  'Technical Skills & Project Tracks'),
        ('7',  'Professional Skills + Supervision'),
        ('8',  'Implementation Tools + Compliance'),
    ]
    right_items = [
        ('9',  'Scenario: Cybersecurity'),
        ('10', ''),
        ('11', 'Tabletop Scenario: Data Insights Dashboard'),
        ('12', ''),
        ('13', 'Scenario: Workflow Automation Prototype'),
        ('14', ''),
        ('15', 'Scenario: Workplace Problem Solving & App Concept Design'),
        ('16', ''),
    ]

    row_h = Inches(0.38)
    for i, (num, label) in enumerate(left_items):
        iy = y + i * row_h
        rect(s, ML, iy, Inches(0.32), Inches(0.28), NAVY)
        txt(s, ML + Pt(2), iy + Pt(2), Inches(0.28), Inches(0.24),
            num, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(s, ML + Inches(0.38), iy + Pt(3), C1W - Inches(0.4), Inches(0.3),
            label, size=9, color=DARK)

    right_items_clean = [(n, l) for n, l in right_items if l]
    for i, (num, label) in enumerate(right_items_clean):
        iy = y + i * row_h
        rect(s, C2X, iy, Inches(0.32), Inches(0.28), NAVY)
        txt(s, C2X + Pt(2), iy + Pt(2), Inches(0.28), Inches(0.24),
            num, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(s, C2X + Inches(0.38), iy + Pt(3), C2W - Inches(0.4), Inches(0.3),
            label, size=9, color=DARK)

    print('  Slide 2 done')

# ── Slide 3: Why Host + Portrait of an Intern ─────────────────────────────────

def slide3(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'Why host + portrait of an intern', 3)

    y = BODY_Y

    # Title
    txt(s, ML, y, CW, Inches(0.35),
        'Why host a SYEP tech intern?', size=14, color=NAVY, bold=True)
    y += Inches(0.38)

    # Intro paragraph
    txt(s, ML, y, CW, Inches(0.55),
        'Hosting a SYEP tech intern creates a mutually beneficial experience — interns gain '
        'real-world skills aligned to the Portrait of a Graduate, while employers gain fresh '
        'perspectives, expanded capacity, and a pipeline of diverse future talent.',
        size=8.5, color=GRAY, wrap=True)
    y += Inches(0.6)

    # Left column
    lx = ML
    ly = y

    # Stats badges
    rect(s, lx, ly, Inches(1.1), Inches(0.7), NAVY)
    txt(s, lx + Pt(4), ly + Pt(4), Inches(1.0), Inches(0.3),
        '6', size=22, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, lx + Pt(4), ly + Pt(34), Inches(1.0), Inches(0.25),
        'Weeks of structured\nlearning', size=6.5, color=WHITE, align=PP_ALIGN.CENTER)

    rect(s, lx + Inches(1.2), ly, Inches(1.1), Inches(0.7), NAVY)
    txt(s, lx + Inches(1.2) + Pt(4), ly + Pt(4), Inches(1.0), Inches(0.3),
        '25', size=22, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, lx + Inches(1.2) + Pt(4), ly + Pt(34), Inches(1.0), Inches(0.25),
        'Hours/week\nSYEP tech intern', size=6.5, color=WHITE, align=PP_ALIGN.CENTER)

    ly += Inches(0.82)

    # Benefits section
    label_bar(s, lx, ly, C1W, 'Benefits to Your Organization')
    ly += Inches(0.26)

    benefits = [
        'Fresh perspectives and innovative ideas on real projects',
        'Expanded project capacity at no direct labor cost',
        'Early access to diverse NYC talent pipeline',
        'Strengthen community ties and corporate responsibility',
        'Mentor and support emerging talent',
    ]
    for b in benefits:
        txt(s, lx + Pt(10), ly, C1W - Pt(12), Inches(0.24),
            f'\u2022  {b}', size=8.5, color=DARK, wrap=True)
        ly += Inches(0.26)

    ly += Inches(0.1)
    # Impact statement
    rect(s, lx, ly, C1W, Inches(0.85), STRIPE)
    txt(s, lx + Pt(6), ly + Pt(4), C1W - Pt(8), Inches(0.18),
        'Impact Statement:', size=8.5, color=NAVY, bold=True)
    txt(s, lx + Pt(6), ly + Pt(22), C1W - Pt(8), Inches(0.62),
        'Each intern you host supports a NYC student in developing skills needed to think '
        'critically, communicate effectively, and grow as a future-ready contributor.',
        size=8, color=DARK, wrap=True)

    # Right column — 3 navy panels
    ry = y
    panels = [
        ('Who They Are', [
            'NYC high school students, ages 16–21',
            'Enrolled in SYEP through DYCD-funded programs',
            'Curious, motivated, and eager to learn',
            'Bring varying levels of prior tech experience',
        ]),
        ('What They Bring', [
            'Curiosity and creativity in approaching new challenges',
            'Emerging critical thinking and problem-solving skills',
            'Developing communication skills across diverse environments',
            'Unique perspectives and lived experiences',
        ]),
        ('What They Need', [
            'Clear structure, expectations, and onboarding',
            'A dedicated supervisor/mentor (1:1 weekly check-ins)',
            'Authentic, meaningful project work',
            'Regular feedback and encouragement',
            'Safe, professional, and inclusive environment',
        ]),
    ]
    for title, bullets in panels:
        rect(s, C2X, ry, C2W, Inches(0.26), NAVY)
        txt(s, C2X + Pt(6), ry + Pt(3), C2W - Pt(8), Inches(0.22),
            title, size=8.5, color=WHITE, bold=True)
        ry += Inches(0.26)
        for b in bullets:
            txt(s, C2X + Pt(10), ry, C2W - Pt(12), Inches(0.24),
                f'\u2022  {b}', size=8, color=DARK, wrap=True)
            ry += Inches(0.25)
        ry += Inches(0.12)

    print('  Slide 3 done')

# ── Slide 4: 6-Week Journey ───────────────────────────────────────────────────

def slide4(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'The 6-week internship journey', 4)

    y = BODY_Y
    txt(s, ML, y, CW, Inches(0.35),
        'The 6-week journey', size=14, color=NAVY, bold=True)
    y += Inches(0.4)

    # Supervisor tip
    rect(s, ML, y, CW, Inches(0.4), STRIPE)
    txt(s, ML + Pt(8), y + Pt(5), CW - Pt(10), Inches(0.32),
        '\U0001f4a1 Supervisor Tip: Each phase builds on the last. Interns who have a strong Explore phase '
        'are better prepared to Build confidently. Don\'t rush onboarding — it sets the foundation for everything.',
        size=8, color=NAVY, wrap=True)
    y += Inches(0.46)

    phases = [
        ('Phase 1 — Explore', 'Weeks 1–2', [
            'Onboarding & company overview',
            'Meet the team & org culture',
            'Orientation to tools and systems',
            'Introduction to internship project',
            'Set goals & learning objectives',
            'Establish norms and expectations',
            'Tour facilities / shadow staff',
            'Complete DYCD compliance docs',
        ]),
        ('Phase 2 — Build', 'Weeks 3–4', [
            'Deep dive into project work',
            'Research, data collection, or prototyping',
            'Apply technical skills to real problem',
            'Collaborative team work sessions',
            'Receive and incorporate feedback',
            'Mid-program check-in with supervisor',
            'Skill-building workshops (as available)',
            'Document findings, tasks, and process',
        ]),
        ('Phase 3 — Apply', 'Week 5', [
            'Finalize project deliverables',
            'Prepare final presentation',
            'Practice presenting to peers',
            'Refine based on feedback',
            'Iterate on tasks using feedback',
        ]),
        ('Phase 4 — Reflect', 'Week 5 (cont.)', [
            'Personal reflection journal',
            'What did I learn? What challenged me?',
            'How have my skills grown?',
            'What would I do differently?',
            'Reflection on tasks completed for the organization',
        ]),
        ('Phase 5 — Launch', 'Week 6', [
            'Career pathway exploration',
            'Identify mentors & next steps',
            'Final presentation to employer leadership',
            'Demo or showcase deliverables',
            'Receive formal evaluation + feedback',
            'Celebrate accomplishments',
            'Receive letter of completion / reference',
            'Network with professionals',
            'Complete exit survey (DYCD)',
        ]),
    ]

    # Phase grid — 5 across, each ~1.44" wide
    ph_w = CW / 5
    ph_x = ML
    ph_h_hdr = Inches(0.5)

    for i, (phase, weeks, items) in enumerate(phases):
        px = ph_x + i * ph_w
        py = y
        # Phase header
        rect(s, px, py, ph_w - Inches(0.04), ph_h_hdr, NAVY)
        txt(s, px + Pt(4), py + Pt(3), ph_w - Pt(8), Inches(0.22),
            phase, size=7.5, color=WHITE, bold=True, wrap=True)
        txt(s, px + Pt(4), py + Pt(28), ph_w - Pt(8), Inches(0.18),
            weeks, size=7, color=ORANGE)

        # Bullets
        by = py + ph_h_hdr + Pt(3)
        for item in items:
            txt(s, px + Pt(4), by, ph_w - Pt(8), Inches(0.2),
                f'\u2022 {item}', size=7, color=DARK, wrap=True)
            by += Inches(0.22)

    # Key deliverables section
    del_y = y + Inches(3.05)
    label_bar(s, ML, del_y, CW, 'Key deliverables', fill=ORANGE)
    del_y += Inches(0.26)

    deliverables = [
        'Weekly 1:1 check-in notes',
        'Timesheets (DYCD compliance)',
        'Mid-program reflection',
        'Final project + presentation',
        'Completed evaluation rubric',
        'Exit survey + certificate',
    ]
    # Two rows of 3
    dw = CW / 3
    for i, d in enumerate(deliverables):
        dx = ML + (i % 3) * dw
        dy = del_y + (i // 3) * Inches(0.26)
        rect(s, dx, dy, dw - Pt(4), Inches(0.22), STRIPE)
        txt(s, dx + Pt(6), dy + Pt(3), dw - Pt(10), Inches(0.2),
            f'\u2713  {d}', size=8, color=DARK)

    print('  Slide 4 done')

# ── Slide 5: Weekly Schedule ──────────────────────────────────────────────────

def slide5(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'Weekly schedule & 25-hour breakdown', 5)

    y = BODY_Y
    txt(s, ML, y, CW, Inches(0.32),
        'Weekly schedule', size=13, color=NAVY, bold=True)
    y += Inches(0.3)
    txt(s, ML, y, CW, Inches(0.3),
        'Interns work 25 hours per week. Below is the recommended daily structure across the 6-week program.',
        size=8.5, color=GRAY)
    y += Inches(0.34)

    # Schedule table
    sched_headers = ['Time block', 'Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday']
    sched_rows = [
        ['9:00–9:30AM\n30 min',
         'Morning Check-In / Stand-up', 'Morning Check-In / Stand-up',
         'Morning Check-In / Stand-up', 'Morning Check-In / Stand-up',
         'Morning Check-In / Stand-up'],
        ['9:30–11:30AM\n2 hrs',
         'Skill Building / Training', 'Project Work Block',
         'Collaborative Workshop', 'Project Work Block',
         'Team Presentation Practice'],
        ['11:30AM–12:30PM\n1 hr',
         'Lunch Break', 'Lunch Break', 'Lunch Break', 'Lunch Break', 'Lunch Break'],
        ['12:30–2:30PM\n2 hrs',
         'Project Work / Research', '1:1 Check-In (rotate)',
         'Guest Speaker / Site Tour', 'Project Work / Research',
         'Weekly Wrap-Up + Reflection'],
        ['2:30–3:00PM\n30 min',
         'Documentation / Journal', 'Documentation / Journal',
         'Documentation / Journal', 'Documentation / Journal',
         'Goal-Setting for Next Week'],
    ]
    col_w = [Inches(1.1), Inches(1.28), Inches(1.28), Inches(1.28), Inches(1.28), Inches(1.28)]
    make_table(s, ML, y, CW, Inches(2.1),
               sched_headers, sched_rows, col_widths=col_w,
               row_size=7.5, hdr_size=8)
    y += Inches(2.18)

    # Note
    txt(s, ML, y, CW, Inches(0.45),
        'Note: Time allocations are a general guide. In the Explore phase (Weeks 1–2), interns may benefit from '
        'more learning and observation. In later phases, time shifts toward project execution and presentation.',
        size=7.5, color=GRAY, italic=True, wrap=True)
    y += Inches(0.5)

    # Hours breakdown
    txt(s, ML, y, Inches(3.6), Inches(0.28),
        'Weekly hours breakdown', size=11, color=NAVY, bold=True)
    txt(s, C2X, y, C2W, Inches(0.28),
        'Compliance requirements', size=11, color=NAVY, bold=True)
    y += Inches(0.3)

    # Hours table
    hrs_headers = ['Activity category', 'Hours/week', '% of time']
    hrs_rows = [
        ['Project Work (core deliverables)', '10 hrs', '40%'],
        ['Skill Building & Training', '5 hrs', '20%'],
        ['Team Collaboration & Meetings', '4 hrs', '16%'],
        ['Professional Development', '3 hrs', '12%'],
        ['Reflection & Documentation', '2 hrs', '8%'],
        ['1:1 Check-In & Feedback', '1 hr', '4%'],
    ]
    hw = [Inches(2.3), Inches(0.7), Inches(0.7)]
    make_table(s, ML, y, Inches(3.7), Inches(1.8),
               hrs_headers, hrs_rows, col_widths=hw, row_size=8)

    # Compliance notes (right)
    cx = C2X
    cy = y
    rect(s, cx, cy, C2W, Inches(0.26), ORANGE)
    txt(s, cx + Pt(6), cy + Pt(3), C2W - Pt(8), Inches(0.22),
        '6-week hour total: 150 hrs', size=9, color=WHITE, bold=True)
    cy += Inches(0.28)

    comp = [
        ('\U0001f4cb Timesheets must be signed by the supervisor weekly and submitted '
         'to the program coordinator by end of business Friday.', STRIPE),
        ('Absences must be reported same day', WHITE),
        ('Maximum 3 absences before review', STRIPE),
        ('Make-up hours require supervisor approval', WHITE),
        ('DYCD site visits may occur unannounced', STRIPE),
    ]
    for note, bg in comp:
        rect(s, cx, cy, C2W, Inches(0.26), bg)
        txt(s, cx + Pt(6), cy + Pt(3), C2W - Pt(8), Inches(0.23),
            note, size=7.5, color=DARK, wrap=True)
        cy += Inches(0.27)

    rect(s, cx, cy, C2W, Inches(0.26), NAVY)
    txt(s, cx + Pt(6), cy + Pt(3), C2W - Pt(8), Inches(0.22),
        'Key Requirement: All interns must maintain 25 hours/week. '
        'Attendance and punctuality are tracked by DYCD as compliance requirements.',
        size=7.5, color=WHITE, wrap=True)

    print('  Slide 5 done')

# ── Slide 6: Technical Skills ─────────────────────────────────────────────────

def slide6(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'Technical skills & project tracks', 6)

    y = BODY_Y
    txt(s, ML, y, CW, Inches(0.32),
        'Technical skills framework', size=13, color=NAVY, bold=True)
    y += Inches(0.3)
    txt(s, ML, y, CW, Inches(0.3),
        'Interns are introduced to foundational and applied technical skills across four core areas.',
        size=8.5, color=GRAY)
    y += Inches(0.34)

    # 4 skill boxes — left column (stacked)
    skill_boxes = [
        ('Cybersecurity', [
            'Threat analysis fundamentals',
            'Login/access data review',
            'Incident documentation',
            'Risk identification',
            'Data privacy concepts',
            'Tabletop exercise facilitation',
        ], 'Example project: Interns may analyze a simulated dataset for suspicious access patterns, '
           'document findings, and discuss potential incident response actions.'),
        ('Data analytics', [
            'Spreadsheet analysis (Excel/Sheets)',
            'Data cleaning & organization',
            'Chart and graph creation',
            'Dashboard building (Tableau/Power BI)',
            'Data storytelling',
            'Insight presentation',
        ], 'Example project: Interns may clean, analyze, and visualize a business dataset to identify '
           'trends and communicate insights through a simple dashboard.'),
        ('Workflow automation', [
            'Process mapping',
            'Zapier / Make.com basics',
            'Google Apps Script',
            'Automation trigger design',
            'Testing & debugging',
            'Efficiency measurement',
        ], 'Example project: Interns may map a manual process, identify inefficiencies, and prototype '
           'a simple automation using accessible tools.'),
        ('Web & coding', [
            'HTML/CSS fundamentals',
            'Basic Python / JavaScript',
            'GitHub & version control',
            'No-code tool exploration',
            'UI/UX prototyping',
            'Technical writing',
        ], 'Example project: Interns may identify a real workplace challenge and design a technology-enabled solution.'),
    ]

    box_h = Inches(1.55)
    for i, (title, items, example) in enumerate(skill_boxes):
        by = y + i * (box_h + Inches(0.06))
        # Title bar
        rect(s, ML, by, C1W, Inches(0.26), NAVY)
        txt(s, ML + Pt(6), by + Pt(3), C1W - Pt(8), Inches(0.22),
            title, size=9, color=WHITE, bold=True)
        # Bullet items — two sub-columns
        half = len(items) // 2
        for j, item in enumerate(items[:half]):
            txt(s, ML + Pt(8), by + Inches(0.3) + j * Inches(0.2),
                C1W / 2 - Pt(10), Inches(0.2),
                f'\u2022 {item}', size=7.5, color=DARK, wrap=True)
        for j, item in enumerate(items[half:]):
            txt(s, ML + C1W / 2, by + Inches(0.3) + j * Inches(0.2),
                C1W / 2 - Pt(4), Inches(0.2),
                f'\u2022 {item}', size=7.5, color=DARK, wrap=True)

    # Right column: Example Project Tracks
    label_bar(s, C2X, y, C2W, 'Example Project Tracks & Applications')
    ry = y + Inches(0.26)

    tracks = [
        ('Track A — Cybersecurity Investigation',
         'Interns may analyze a simulated dataset for suspicious access patterns, document findings, '
         'and discuss potential incident response actions.'),
        ('Track B — Data Insights Dashboard',
         'Interns may clean, analyze, and visualize a business dataset to identify trends and '
         'communicate insights through a simple dashboard.'),
        ('Track C — Workflow Automation Prototype',
         'Interns may map a manual process, identify inefficiencies, and prototype a simple '
         'automation using accessible tools.'),
        ('Track D — Workplace Problem Solving & App Concept Design',
         'Interns may identify a real workplace challenge, conduct research with staff or customers, '
         'and design a technology-enabled solution (app, tool, or system) to improve how work gets done.'),
    ]

    track_fills = [STRIPE, WHITE, STRIPE, WHITE]
    for i, (track_title, desc) in enumerate(tracks):
        th = Inches(1.52)
        rect(s, C2X, ry, C2W, th, track_fills[i])
        txt(s, C2X + Pt(6), ry + Pt(4), C2W - Pt(8), Inches(0.22),
            track_title, size=8, color=NAVY, bold=True, wrap=True)
        txt(s, C2X + Pt(8), ry + Pt(26), C2W - Pt(10), th - Pt(28),
            desc, size=7.5, color=DARK, wrap=True)

        # Project Selection note only after Track D
        if i == 3:
            txt(s, C2X + Pt(6), ry + Pt(74), C2W - Pt(8), Inches(0.3),
                'Project Selection: Employers are encouraged to design projects aligned to their '
                "organization's work, using these examples as a starting point.",
                size=7, color=GRAY, italic=True, wrap=True)

        ry += th + Inches(0.04)

    print('  Slide 6 done')

# ── Slide 7: Professional Skills ─────────────────────────────────────────────

def slide7(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'Professional skills + supervision best practices', 7)

    y = BODY_Y
    txt(s, ML, y, CW, Inches(0.32),
        'Professional skills + supervision best practices', size=13, color=NAVY, bold=True)
    y += Inches(0.36)
    txt(s, ML, y, CW, Inches(0.3),
        'Beyond technical skills, SYEP interns develop essential professional competencies. '
        'Supervisors play a key role in modeling and reinforcing these behaviors.',
        size=8.5, color=GRAY, wrap=True)
    y += Inches(0.38)

    # Left: Competency table
    txt(s, ML, y, C1W, Inches(0.26),
        'Workplace & professional skills', size=10, color=NAVY, bold=True)
    y += Inches(0.3)

    comp_headers = ['Competency', 'How interns develop this']
    comp_rows = [
        ['Effective Communicator',
         'Stand-ups, presentations, written updates, email norms'],
        ['Global Citizen',
         'Team roles, collaboration, group problem-solving, etiquette'],
        ['Critical Thinker',
         'Project analysis, peer review, retrospectives, research'],
        ['Creative Innovator',
         'Adapting scope, prototyping, exploring new approaches'],
        ['Reflective / Future Focused',
         'Goal-setting, time tracking, documentation, reflection'],
    ]
    cw1 = [Inches(1.4), Inches(1.9)]
    make_table(s, ML, y, C1W, Inches(2.2), comp_headers, comp_rows,
               col_widths=cw1, row_size=8)
    y += Inches(2.28)

    # Career exploration note
    rect(s, ML, y, C1W, Inches(0.6), STRIPE)
    txt(s, ML + Pt(6), y + Pt(4), C1W - Pt(8), Inches(0.55),
        '\U0001f4cc Career Exploration: Dedicate at least 30 minutes per week for interns to reflect '
        'on their tasks and explore related career pathways, job titles, and salary data.',
        size=8, color=NAVY, wrap=True)
    y += Inches(0.68)

    # Right: Supervision best practices
    ry = BODY_Y + Inches(0.66)
    txt(s, C2X, ry, C2W, Inches(0.26),
        'Supervision best practices', size=10, color=NAVY, bold=True)
    ry += Inches(0.3)

    label_bar(s, C2X, ry, C2W, 'Before the Internship Starts', fill=LNAVY, h=Inches(0.22))
    ry += Inches(0.24)
    before = [
        'Identify a backup supervisor in case of absence',
        'Brief your team on hosting a high school intern',
        'Review DYCD guidelines and compliance requirements',
        'Prepare a Week 1 onboarding plan — see resources',
        'Set up workstation, accounts, and tools for day-to-day tasks',
    ]
    for b in before:
        txt(s, C2X + Pt(8), ry, C2W - Pt(10), Inches(0.22),
            f'\u2022  {b}', size=8, color=DARK, wrap=True)
        ry += Inches(0.24)

    ry += Inches(0.08)
    label_bar(s, C2X, ry, C2W, 'During the Internship', fill=LNAVY, h=Inches(0.22))
    ry += Inches(0.24)
    during = [
        'Conduct weekly 1:1 check-ins (use template)',
        'Assign interns meaningful day-to-day tasks — not busy work',
        'Introduce interns to professionals and team workflows',
        'Provide specific, actionable feedback regularly',
        'Track attendance and sign timesheets weekly',
        'Contact program coordinator if issues arise',
    ]
    for d in during:
        txt(s, C2X + Pt(8), ry, C2W - Pt(10), Inches(0.22),
            f'\u2022  {d}', size=8, color=DARK, wrap=True)
        ry += Inches(0.24)

    ry += Inches(0.08)
    rect(s, C2X, ry, C2W, Inches(0.22), STRIPE)
    txt(s, C2X + Pt(6), ry + Pt(3), C2W - Pt(8), Inches(0.2),
        'Effective Feedback Framework', size=8, color=NAVY, bold=True)
    ry += Inches(0.24)
    feedback = [
        ('Strengths first:', 'Start with what\'s working well'),
        ('Be specific:', '"Your chart was clear because..." not "good job"'),
        ('Growth-oriented:', 'Frame challenges from assigned tasks as opportunities'),
        ('Action-focused:', 'Always include a next step'),
        ('Regular cadence:', "Don't save everything for Week 6"),
    ]
    for label, desc in feedback:
        tb = s.shapes.add_textbox(C2X + Pt(8), ry, C2W - Pt(10), Inches(0.22))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = label + ' '; r1.font.bold = True
        r1.font.size = Pt(8); r1.font.color.rgb = NAVY
        r2 = p.add_run(); r2.text = desc
        r2.font.size = Pt(8); r2.font.color.rgb = DARK
        ry += Inches(0.24)

    print('  Slide 7 done')

# ── Slide 8: Implementation Tools ────────────────────────────────────────────

def slide8(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'Implementation tools + compliance', 8)

    y = BODY_Y
    txt(s, ML, y, CW, Inches(0.32),
        'Implementation tools + compliance', size=13, color=NAVY, bold=True)
    y += Inches(0.36)

    # Left column
    txt(s, ML, y, C1W, Inches(0.26),
        'Recommended digital tools', size=10, color=NAVY, bold=True)
    y += Inches(0.28)

    tools_headers = ['Category', 'Tools']
    tools_rows = [
        ['Productivity', 'Google Workspace (Docs, Sheets, Slides), Microsoft 365'],
        ['Project Management', 'Trello, Notion, Asana (free tiers)'],
        ['Communication', 'Slack, Email, Google Meet / Zoom'],
        ['Data Analysis', 'Excel, Google Sheets, Tableau Public'],
        ['Automation', 'Zapier (free), Make.com, Google Apps Script'],
        ['Design & Presentation', 'Canva, Google Slides, PowerPoint'],
        ['Version Control', 'GitHub (public repos, free)'],
    ]
    cw = [Inches(1.15), Inches(2.2)]
    make_table(s, ML, y, C1W, Inches(2.0), tools_headers, tools_rows,
               col_widths=cw, row_size=8)
    y += Inches(2.08)

    # Required documentation
    label_bar(s, ML, y, C1W, 'Required documentation', fill=LNAVY, h=Inches(0.22))
    y += Inches(0.24)
    docs = [
        'Intern Learning Plan (use the Learning Plan Template)',
        'Weekly 1:1 Check-In Notes (use standalone template)',
        'Attendance Tracking (per SYEP requirements)',
        'Incident Reports (if applicable; submit within 24 hours)',
    ]
    for d in docs:
        txt(s, ML + Pt(8), y, C1W - Pt(10), Inches(0.22),
            f'\u2022  {d}', size=8, color=DARK, wrap=True)
        y += Inches(0.24)

    # Right column
    ry = BODY_Y + Inches(0.36)
    txt(s, C2X, ry, C2W, Inches(0.26),
        'Onboarding checklist', size=10, color=NAVY, bold=True)
    ry += Inches(0.28)

    checklist = [
        'Review and sign DYCD participation agreement',
        'Complete safety orientation and emergency contact form',
        'Receive workstation / device assignment',
        'Set up company email and required accounts',
        'Tour the office and meet team members',
        'Review internship project brief and goals',
        'Sign technology use and confidentiality agreement',
        'Set Week 1 goals with supervisor',
    ]
    for item in checklist:
        rect(s, C2X, ry, C2W, Inches(0.24), STRIPE if checklist.index(item) % 2 else WHITE)
        txt(s, C2X + Pt(6), ry + Pt(2), C2W - Pt(8), Inches(0.22),
            f'\u2610  {item}', size=8, color=DARK, wrap=True)
        ry += Inches(0.25)

    ry += Inches(0.08)
    # Compliance requirements
    rect(s, C2X, ry, C2W, Inches(0.22), ORANGE)
    txt(s, C2X + Pt(6), ry + Pt(3), C2W - Pt(8), Inches(0.2),
        '\u26a0\ufe0f DYCD Compliance Requirements', size=8.5, color=WHITE, bold=True)
    ry += Inches(0.24)

    compliance = [
        'Interns must not be left unsupervised at any time',
        'No task may put an intern\'s physical safety at risk',
        'Zero tolerance for harassment, bullying, or discrimination',
        'Interns may not handle sensitive client/financial data without oversight',
        'All incidents must be reported to DYCD within 24 hours',
    ]
    for c in compliance:
        txt(s, C2X + Pt(8), ry, C2W - Pt(10), Inches(0.22),
            f'\u26a0  {c}', size=8, color=DARK, wrap=True)
        ry += Inches(0.24)

    txt(s, C2X, ry + Pt(4), C2W, Inches(0.3),
        '\u26a0\ufe0f ALL of the above are required to remain in compliance with DYCD SYEP guidelines. '
        'Non-compliance may result in program removal.',
        size=7.5, color=ORANGE, bold=True, wrap=True)

    print('  Slide 8 done')

# ── Slide 9: Scenario A Intro ─────────────────────────────────────────────────

def slide9(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'NYC DYCD — Scenario A: Cybersecurity Tabletop Investigation', 9,
            'NYC DYCD — Scenario A: Cybersecurity Tabletop Investigation')

    y = BODY_Y

    # Scenario banner
    rect(s, 0, y, W, Inches(0.3), ORANGE)
    txt(s, ML, y + Pt(3), CW, Inches(0.26),
        'PROJECT SCENARIO A  —  TRACK: CYBERSECURITY', size=9, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)
    y += Inches(0.34)

    # Subheader
    txt(s, ML, y, CW, Inches(0.28),
        'Cybersecurity tabletop investigation + friction log', size=12, color=NAVY, bold=True)
    y += Inches(0.3)
    txt(s, ML, y, CW, Inches(0.2),
        'Host Company: NovaBank (Fintech Startup)  |  Duration: 6 Weeks  |  Team Size: 1–5 Interns',
        size=8, color=GRAY)
    y += Inches(0.26)

    # Left column
    lx = ML
    ly = y

    # Company background
    label_bar(s, lx, ly, C1W, 'Company background')
    ly += Inches(0.26)
    txt(s, lx, ly, C1W, Inches(0.85),
        'NovaBank is a growing NYC-based fintech startup providing digital banking services to '
        'underserved communities. Growing usage has raised concerns about unauthorized access and '
        'potential fraud. Your intern team will act as a junior cybersecurity task force, supervised '
        'by employer staff.',
        size=8, color=DARK, wrap=True)
    ly += Inches(0.9)

    # The challenge
    label_bar(s, lx, ly, C1W, 'The challenge')
    ly += Inches(0.26)
    txt(s, lx, ly, C1W, Inches(0.7),
        "NovaBank's security team has flagged an anomaly: a suspicious spike in failed login attempts "
        "over the past 30 days. The team needs to analyze login activity data, identify patterns of "
        "suspicious behavior, document findings, and recommend a response.",
        size=8, color=DARK, wrap=True)
    ly += Inches(0.78)

    # Mission callout
    rect(s, lx, ly, C1W, Inches(0.55), STRIPE)
    txt(s, lx + Pt(6), ly + Pt(4), C1W - Pt(8), Inches(0.5),
        '\u26a1 Intern Mission: Analyze the simulated login dataset, identify indicators of compromise '
        '(IOC), build a friction log documenting suspicious events, and present your findings and '
        'recommendations to NovaBank leadership.',
        size=8, color=NAVY, wrap=True)
    ly += Inches(0.62)

    # Dataset options
    label_bar(s, lx, ly, C1W, 'Dataset options')
    ly += Inches(0.26)
    options = [
        ('Option A — Kaggle Dataset:',
         'Use the "Login Attempts" dataset from Kaggle.com (free, publicly available). '
         'Contains 10,000+ rows of login data.'),
        ('Option B — Anonymized Internal Data:',
         'NovaBank will provide a sanitized, anonymized export from their test environment.'),
        ('Option C — Simulated Dataset:',
         'Program coordinator provides a purpose-built 500-row dataset pre-loaded with anomalies.'),
    ]
    for label, desc in options:
        tb = s.shapes.add_textbox(lx + Pt(6), ly, C1W - Pt(8), Inches(0.32))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = label + ' '; r1.font.bold = True
        r1.font.size = Pt(7.5); r1.font.color.rgb = NAVY
        r2 = p.add_run(); r2.text = desc
        r2.font.size = Pt(7.5); r2.font.color.rgb = DARK
        ly += Inches(0.32)

    # Right column
    rx = C2X
    ry = y

    # Team roles
    label_bar(s, rx, ry, C2W, 'Team roles')
    ry += Inches(0.26)

    roles_headers = ['Role', 'Responsibilities']
    roles_rows = [
        [{'text': 'Security Analyst', 'bold': False},
         'Reviews login data to identify suspicious activity, flags anomalies, and documents potential security risks'],
        [{'text': 'Data Analyst', 'bold': True, 'color': NAVY},
         'Cleans and organizes dataset; builds charts showing login trends over time'],
        ['Research Lead',
         'Researches common cyberattack types (brute force, credential stuffing); documents background'],
        ['Documentation Lead',
         'Maintains friction log; records all flagged events with timestamps and notes'],
        ['Presentation Lead',
         'Compiles team findings; creates and presents the final slide deck'],
    ]
    rcw = [Inches(1.2), Inches(2.52)]
    make_table(s, rx, ry, C2W, Inches(1.85), roles_headers, roles_rows,
               col_widths=rcw, row_size=7.5)
    ry += Inches(1.92)

    # Friction log
    label_bar(s, rx, ry, C2W, 'Friction log template')
    ry += Inches(0.26)

    fl_headers = ['#', 'Timestamp', 'User', 'IP address', 'Event type', 'Severity', 'Notes']
    fl_rows = [
        ['1', 'MM/DD\nHH:MM', 'USR-####', 'xxx.xxx.x.x', 'Failed Login', '\u26a0\ufe0f Medium', '3x in 2 min'],
        ['2', '', '', '', '', '', ''],
        ['3', '', '', '', '', '', ''],
    ]
    flcw = [Inches(0.22), Inches(0.55), Inches(0.55), Inches(0.68), Inches(0.62), Inches(0.6), Inches(0.5)]
    make_table(s, rx, ry, C2W, Inches(0.8), fl_headers, fl_rows,
               col_widths=flcw, row_size=7)
    ry += Inches(0.88)

    txt(s, rx, ry, C2W, Inches(0.3),
        'Severity: \U0001f534 Critical = 10+ failed/5 min  |  \u26a0\ufe0f Medium = unusual IP  |  \U0001f7e1 Low = off-hours',
        size=7, color=GRAY, wrap=True)

    print('  Slide 9 done')

# ── Slide 10: Scenario A — 6-Week Structure ───────────────────────────────────

def slide10(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'NYC DYCD — Scenario A: Cybersecurity Tabletop Investigation', 10,
            'NYC DYCD — Scenario A: Cybersecurity Tabletop Investigation')

    y = BODY_Y
    txt(s, ML, y, CW, Inches(0.32),
        'Scenario A: Cybersecurity — 6-week structure & deliverables', size=12, color=NAVY, bold=True)
    y += Inches(0.35)

    label_bar(s, ML, y, CW, '6-week project structure', fill=ORANGE)
    y += Inches(0.27)

    wk_headers = ['Week', 'Phase', 'Activities', 'Deliverable']
    wk_rows = [
        ['Week 1', 'Explore',
         'Onboard to NovaBank; learn about digital banking security; watch intro to cybersecurity videos; '
         'meet team; review DYCD expectations; explore dataset structure',
         'Team roles assigned; dataset downloaded; intro to cybersecurity 1-page summary'],
        ['Week 2', 'Explore + Build',
         'Deep-read the dataset; identify columns and data types; research brute force attacks and '
         'credential stuffing; begin flagging obvious anomalies (100+ failed attempts, single IP); '
         'create initial filters in spreadsheet',
         'Annotated dataset with initial flags; research notes on attack types'],
        ['Week 3', 'Build',
         'Systematic anomaly detection: sort by IP, User ID, timestamp; identify top 10 suspicious '
         'accounts; build login trend charts (daily failed logins over 30 days); begin friction log '
         'with Documentation Lead',
         'Top 10 suspicious accounts report; login trend chart; friction log (initial draft)'],
        ['Week 4', 'Build + Apply',
         'Complete friction log with all flagged events; research incident response best practices; '
         'draft 3–5 security recommendations; connect data findings to real-world breach scenarios; '
         'mid-program check-in with supervisor',
         'Complete friction log; written recommendations draft; mid-program self-assessment'],
        ['Week 5', 'Apply',
         'Synthesize findings into presentation narrative; create visuals for final slide deck (charts, '
         'friction log summary, severity heatmap); practice presenting; incorporate supervisor feedback; '
         'write executive summary (1-page)',
         'Final slide deck draft; executive summary; practice presentation completed'],
        ['Week 6', 'Launch + Reflect',
         'Final presentation to NovaBank leadership (15–20 min); Q&A with leadership panel; receive '
         'formal evaluation; team celebration; complete DYCD exit survey; write personal reflection',
         'Final presentation delivered; evaluation form signed; exit survey completed; personal reflection'],
    ]
    col_w = [Inches(0.62), Inches(0.85), Inches(3.8), Inches(2.13)]
    make_table(s, ML, y, CW, Inches(4.2), wk_headers, wk_rows,
               col_widths=col_w, row_size=7.5)
    y += Inches(4.28)

    # Expected deliverables + learning outcomes
    txt(s, ML, y, C1W, Inches(0.24),
        'Expected final deliverables', size=9, color=NAVY, bold=True)
    txt(s, C2X, y, C2W, Inches(0.24),
        'Key learning outcomes', size=9, color=NAVY, bold=True)
    y += Inches(0.26)

    deliverables = [
        '\u2705 Complete Friction Log — all suspicious events documented with timestamps, severity, and notes',
        '\u2705 Anomaly Analysis Report — top 10 flagged accounts with evidence and reasoning',
        '\u2705 Login Trend Visualizations — charts showing patterns over 30-day window',
        '\u2705 3–5 Security Recommendations — specific, actionable steps for NovaBank',
        '\u2705 Final Presentation Deck — 10–15 slides, presented to leadership',
        '\u2705 1-Page Executive Summary — written for a non-technical audience',
    ]
    for d in deliverables:
        txt(s, ML, y, C1W, Inches(0.22), d, size=7.5, color=DARK, wrap=True)
        y += Inches(0.23)

    outcomes = [
        'Data analysis', 'Threat identification', 'Documentation',
        'Risk assessment', 'Presentation skills', 'Cybersecurity concepts',
        'Teamwork', 'Critical thinking',
    ]
    oy = y - len(deliverables) * Inches(0.23)
    for o in outcomes:
        rect(s, C2X, oy, C2W, Inches(0.22), STRIPE if outcomes.index(o) % 2 else WHITE)
        txt(s, C2X + Pt(6), oy + Pt(2), C2W - Pt(8), Inches(0.2),
            f'\u2022  {o}', size=8, color=DARK)
        oy += Inches(0.24)

    # Supervisor note
    sy = max(y, oy) + Inches(0.08)
    rect(s, ML, sy, CW, Inches(0.35), STRIPE)
    txt(s, ML + Pt(6), sy + Pt(4), CW - Pt(8), Inches(0.3),
        '\u26a1 Supervisor Guidance: Interns are NOT expected to access real systems. All work uses the '
        'provided/downloaded dataset only. Emphasize ethical boundaries and data privacy throughout.',
        size=7.5, color=NAVY, wrap=True)

    print('  Slide 10 done')

# ── Slide 11: Scenario B Intro ────────────────────────────────────────────────

def slide11(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'NYC DYCD — Scenario B: Data Insights Dashboard', 11,
            'NYC DYCD — Scenario B: Data Insights Dashboard')

    y = BODY_Y
    rect(s, 0, y, W, Inches(0.3), ORANGE)
    txt(s, ML, y + Pt(3), CW, Inches(0.26),
        'PROJECT SCENARIO B  —  TRACK: DATA ANALYTICS', size=9, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)
    y += Inches(0.34)

    txt(s, ML, y, CW, Inches(0.28),
        'Data insights dashboard', size=12, color=NAVY, bold=True)
    y += Inches(0.28)
    txt(s, ML, y, CW, Inches(0.2),
        'Host Company: GreenCart (E-Commerce Grocery Startup)  |  Duration: 6 Weeks  |  Team Size: 1–5 Interns',
        size=8, color=GRAY)
    y += Inches(0.26)

    lx = ML
    ly = y

    label_bar(s, lx, ly, C1W, 'Company background')
    ly += Inches(0.26)
    txt(s, lx, ly, C1W, Inches(0.75),
        'GreenCart is a fast-growing NYC-based e-commerce grocery startup focused on providing fresh, '
        'affordable produce to underserved neighborhoods. Leadership needs better visibility into order '
        'trends, delivery performance, and customer satisfaction to make data-driven decisions.',
        size=8, color=DARK, wrap=True)
    ly += Inches(0.82)

    label_bar(s, lx, ly, C1W, 'The challenge')
    ly += Inches(0.26)
    txt(s, lx, ly, C1W, Inches(0.65),
        "GreenCart's operations team currently tracks orders in a spreadsheet but has no visual dashboard. "
        "Your intern team will analyze their order data and build a visual dashboard to surface actionable insights.",
        size=8, color=DARK, wrap=True)
    ly += Inches(0.72)

    rect(s, lx, ly, C1W, Inches(0.5), STRIPE)
    txt(s, lx + Pt(6), ly + Pt(4), C1W - Pt(8), Inches(0.45),
        '\u26a1 Intern Mission: Clean and analyze GreenCart\'s order/delivery dataset (100–500 rows), '
        'identify 3–5 key business insights, and build an interactive visual dashboard that leadership '
        'can use to make better decisions.',
        size=8, color=NAVY, wrap=True)
    ly += Inches(0.58)

    label_bar(s, lx, ly, C1W, 'Recommended tools')
    ly += Inches(0.26)
    tools = ['Microsoft Excel', 'Google Sheets', 'Tableau Public', 'Power BI (free)', 'Canva (charts)']
    for t in tools:
        txt(s, lx + Pt(8), ly, C1W - Pt(10), Inches(0.2),
            f'\u2022  {t}', size=8, color=DARK)
        ly += Inches(0.21)
    txt(s, lx, ly + Pt(4), C1W, Inches(0.2),
        'Supervisors: Ensure interns have access to at least one of these tools before Week 1 begins.',
        size=7, color=GRAY, italic=True, wrap=True)

    # Right column
    rx = C2X
    ry = y

    label_bar(s, rx, ry, C2W, 'Team roles')
    ry += Inches(0.26)

    roles_rows = [
        ['Data Analyst', 'Cleans dataset; removes duplicates and nulls; standardizes formatting; runs pivot tables'],
        ['Visualization Lead', 'Creates charts and graphs; builds the visual dashboard in chosen tool (Excel/Tableau/Power BI)'],
        ['Insights Analyst', 'Interprets charts; writes up 3–5 key findings in plain language; identifies trends and outliers'],
        ['Research Lead', 'Benchmarks GreenCart\'s metrics against industry standards; researches e-commerce best practices'],
        ['Presentation Lead', 'Compiles all findings; builds final presentation deck; leads the final delivery to GreenCart leadership'],
    ]
    rcw = [Inches(1.3), Inches(2.42)]
    make_table(s, rx, ry, C2W, Inches(1.85), ['Role', 'Responsibilities'], roles_rows,
               col_widths=rcw, row_size=7.5)
    ry += Inches(1.92)

    label_bar(s, rx, ry, C2W, 'The dataset — GreenCart_Orders_Q2.csv')
    ry += Inches(0.26)

    ds_headers = ['Column', 'Description', 'Example']
    ds_rows = [
        ['Order ID', 'Unique order identifier', 'GC-2024-00234'],
        ['Product Category', 'Type of product ordered', 'Vegetables, Fruits, Dairy'],
        ['Order Date', 'Date order was placed', '2024-04-15'],
        ['Delivery Time (hrs)', 'Hours from order to delivery', '3.5'],
        ['City/Borough', "Customer's borough", 'Bronx, Brooklyn, Queens'],
        ['Order Status', 'Completion status', 'Delivered, Returned, Cancelled'],
        ['Customer Rating', '1–5 star rating from customer', '4'],
        ['Order Value ($)', 'Total order amount', '$47.50'],
    ]
    dscw = [Inches(1.2), Inches(1.55), Inches(0.97)]
    make_table(s, rx, ry, C2W, Inches(2.1), ds_headers, ds_rows,
               col_widths=dscw, row_size=7)
    print('  Slide 11 done')

# ── Slide 12: Scenario B — 6-Week Structure ───────────────────────────────────

def slide12(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'NYC DYCD — Scenario B: Data Insights Dashboard', 12,
            'NYC DYCD — Scenario B: Data Insights Dashboard')

    y = BODY_Y
    txt(s, ML, y, CW, Inches(0.32),
        'Scenario B: Data insights dashboard — 6-week structure & deliverables',
        size=12, color=NAVY, bold=True)
    y += Inches(0.35)

    label_bar(s, ML, y, CW, '6-week project structure', fill=ORANGE)
    y += Inches(0.27)

    wk_rows = [
        ['Week 1', 'Explore',
         'Onboard to GreenCart; understand e-commerce business model; explore the dataset structure; '
         'identify what each column means; ask questions about the business; discuss what leadership '
         'wants to know; assign team roles',
         'Dataset exploration notes; 5 business questions the team wants to answer'],
        ['Week 2', 'Explore + Build',
         'Data cleaning: remove blanks, fix date formats, remove duplicate Order IDs; standardize '
         'borough names and product categories; document cleaning steps; run basic statistics; '
         'identify any data quality issues',
         'Cleaned dataset (v1); data cleaning log; basic statistics summary'],
        ['Week 3', 'Build',
         'Build pivot tables and initial charts: orders by borough, top product categories, delivery '
         'time trends, order status breakdown; identify which borough has lowest satisfaction; find '
         'peak order days/times; start dashboard layout design',
         '5+ charts created; dashboard wireframe/sketch; initial insights notes'],
        ['Week 4', 'Build + Apply',
         'Complete interactive dashboard; add filters (by date, borough, category); finalize 3–5 key '
         'insights with supporting data; write "so what" interpretation for each insight; mid-program '
         'check-in; draft recommendations for GreenCart',
         'Dashboard v1 (complete); written insights; draft recommendations; mid-program self-assessment'],
        ['Week 5', 'Apply',
         'Build final presentation deck; use dashboard screenshots as visual evidence; practice '
         'presenting insights in non-technical language; present to peers; incorporate supervisor '
         'feedback; refine dashboard based on feedback',
         'Final presentation deck; dashboard v2 (refined); practice presentation completed'],
        ['Week 6', 'Launch + Reflect',
         'Final presentation to GreenCart leadership (15–20 min); live demo of dashboard; answer '
         'leadership questions; receive formal evaluation; complete DYCD exit survey; write personal '
         'reflection; team celebration',
         'Final presentation delivered; evaluation signed; exit survey; personal reflection'],
    ]
    col_w = [Inches(0.62), Inches(0.85), Inches(3.8), Inches(2.13)]
    make_table(s, ML, y, CW, Inches(4.2), ['Week', 'Phase', 'Activities', 'Deliverable'],
               wk_rows, col_widths=col_w, row_size=7.5)
    y += Inches(4.28)

    txt(s, ML, y, C1W, Inches(0.24),
        'Expected final deliverables', size=9, color=NAVY, bold=True)
    txt(s, C2X, y, C2W, Inches(0.24),
        'Key learning outcomes', size=9, color=NAVY, bold=True)
    y += Inches(0.26)

    deliverables = [
        '\u2705 Clean Dataset — documented cleaning process, no errors/blanks',
        '\u2705 Visual Dashboard — 5+ charts/visualizations, filterable by borough & date',
        '\u2705 3–5 Key Business Insights — written in plain language with data evidence',
        '\u2705 Business Recommendations — actionable suggestions based on findings',
        '\u2705 Final Presentation Deck — 10–15 slides, presented to leadership',
        '\u2705 Industry Benchmarking — how does GreenCart compare to e-commerce averages?',
    ]
    for d in deliverables:
        txt(s, ML, y, C1W, Inches(0.22), d, size=7.5, color=DARK, wrap=True)
        y += Inches(0.23)

    outcomes = ['Data cleaning', 'Excel / Sheets', 'Data visualization',
                'Business analysis', 'Storytelling with data', 'Dashboard design', 'Presentation skills']
    oy = y - len(deliverables) * Inches(0.23)
    for o in outcomes:
        rect(s, C2X, oy, C2W, Inches(0.22), STRIPE if outcomes.index(o) % 2 else WHITE)
        txt(s, C2X + Pt(6), oy + Pt(2), C2W - Pt(8), Inches(0.2),
            f'\u2022  {o}', size=8, color=DARK)
        oy += Inches(0.24)

    sy = max(y, oy) + Inches(0.08)
    rect(s, ML, sy, CW, Inches(0.4), STRIPE)
    txt(s, ML + Pt(6), sy + Pt(4), CW - Pt(8), Inches(0.35),
        '\u26a1 Sample Insight Example: "Brooklyn has the highest order volume but the lowest customer '
        'rating (avg 2.8/5). Delivery times in Brooklyn average 5.2 hrs vs. 2.8 hrs in Queens. '
        'Recommendation: Investigate delivery routing in Brooklyn — this is your biggest opportunity '
        'for satisfaction improvement."',
        size=7.5, color=NAVY, wrap=True)

    print('  Slide 12 done')

# ── Slide 13: Scenario C Intro ────────────────────────────────────────────────

def slide13(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'NYC DYCD — Scenario C: Workflow Automation Prototype', 13,
            'NYC DYCD — Scenario C: Workflow Automation Prototype')

    y = BODY_Y
    rect(s, 0, y, W, Inches(0.3), ORANGE)
    txt(s, ML, y + Pt(3), CW, Inches(0.26),
        'PROJECT SCENARIO C  —  TRACK: WORKFLOW AUTOMATION', size=9, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)
    y += Inches(0.34)

    txt(s, ML, y, CW, Inches(0.28),
        'Workflow automation prototype', size=12, color=NAVY, bold=True)
    y += Inches(0.28)
    txt(s, ML, y, CW, Inches(0.2),
        'Host Company: BrightPath (Tech Consulting Firm)  |  Duration: 6 Weeks  |  Team Size: 1–5 Interns',
        size=8, color=GRAY)
    y += Inches(0.26)

    lx = ML
    ly = y

    label_bar(s, lx, ly, C1W, 'Company background')
    ly += Inches(0.26)
    txt(s, lx, ly, C1W, Inches(0.7),
        'BrightPath is a tech consulting firm that helps small businesses streamline their operations. '
        'Internally, BrightPath still relies on a manual form submission process for new client inquiries '
        '— a time-consuming bottleneck that causes delays and errors. Your intern team will analyze the '
        'current workflow and build an automation prototype to fix it.',
        size=8, color=DARK, wrap=True)
    ly += Inches(0.78)

    label_bar(s, lx, ly, C1W, 'The challenge')
    ly += Inches(0.26)
    txt(s, lx, ly, C1W, Inches(0.65),
        'When a new client submits an inquiry, a staff member must manually copy form data, send a '
        'welcome email, create a task, add the contact to the CRM, and schedule a discovery call. '
        'This 5-step manual process takes 20–30 minutes per inquiry and is prone to errors and delays.',
        size=8, color=DARK, wrap=True)
    ly += Inches(0.72)

    rect(s, lx, ly, C1W, Inches(0.5), STRIPE)
    txt(s, lx + Pt(6), ly + Pt(4), C1W - Pt(8), Inches(0.45),
        '\u26a1 Intern Mission: Map the current manual workflow, identify bottlenecks, design an '
        'automated alternative, build a working prototype using no-code/low-code tools, and measure '
        'time savings.',
        size=8, color=NAVY, wrap=True)
    ly += Inches(0.58)

    label_bar(s, lx, ly, C1W, 'Current manual workflow (5 steps)')
    ly += Inches(0.26)
    steps = [
        'Step 1 — Copy form data to spreadsheet (5 min)',
        'Step 2 — Send welcome email manually (5 min)',
        'Step 3 — Create task in Asana/Trello (5 min)',
        'Step 4 — Add contact to CRM (5 min)',
        'Step 5 — Schedule discovery call (10 min)',
    ]
    for step in steps:
        txt(s, lx + Pt(8), ly, C1W - Pt(10), Inches(0.22),
            f'\u2022  {step}', size=8, color=DARK)
        ly += Inches(0.23)

    # Example outcome
    rect(s, lx, ly, C1W, Inches(0.55), STRIPE)
    txt(s, lx + Pt(6), ly + Pt(4), C1W - Pt(8), Inches(0.5),
        'Example Outcome: Before automation, each inquiry required a 5-step manual process taking '
        '20–30 minutes. After the intern team built the automation prototype, processing time dropped '
        'to under 5 minutes — with fewer errors and improved response consistency.',
        size=7.5, color=DARK, italic=True, wrap=True)

    # Right column
    rx = C2X
    ry = y

    label_bar(s, rx, ry, C2W, 'Team roles')
    ry += Inches(0.26)

    roles_rows = [
        ['Process Analyst', 'Documents current 5-step manual workflow in detail; interviews staff; creates process map'],
        ['Data Analyst', 'Analyzes inquiry dataset to measure current performance (avg response time, error rate, volume)'],
        ['Automation Builder', 'Builds the automation prototype using Zapier, Make.com, or Google Apps Script'],
        ['QA / Tester', 'Tests the prototype with sample data; documents bugs and improvements'],
        ['Presentation Lead', 'Compiles findings; builds final presentation showing before/after comparison'],
    ]
    rcw = [Inches(1.3), Inches(2.42)]
    make_table(s, rx, ry, C2W, Inches(1.65), ['Role', 'Responsibilities'], roles_rows,
               col_widths=rcw, row_size=7.5)
    ry += Inches(1.72)

    label_bar(s, rx, ry, C2W, 'The dataset — BrightPath_Inquiries_Sample.csv')
    ry += Inches(0.26)

    ds_rows = [
        ['Inquiry ID', 'Unique form submission ID'],
        ['Client Name', 'Name of inquiring client'],
        ['Email Address', 'Client contact email'],
        ['Business Type', 'Type of business (retail, services, etc.)'],
        ['Service Interest', 'Which BrightPath service they want'],
        ['Submission Date', 'When the form was submitted'],
        ['Response Date', 'When staff first responded'],
        ['Response Time (hrs)', 'Calculated time to first response'],
        ['Status', 'Pending / In Progress / Closed'],
    ]
    dscw = [Inches(1.55), Inches(2.17)]
    make_table(s, rx, ry, C2W, Inches(2.1), ['Column', 'Description'], ds_rows,
               col_widths=dscw, row_size=7)
    print('  Slide 13 done')

# ── Slide 14: Scenario C — 6-Week Structure ───────────────────────────────────

def slide14(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'NYC DYCD — Scenario C: Workflow Automation Prototype', 14,
            'NYC DYCD — Scenario C: Workflow Automation Prototype')

    y = BODY_Y
    txt(s, ML, y, CW, Inches(0.32),
        'Scenario C: Workflow automation — 6-week structure & deliverables',
        size=12, color=NAVY, bold=True)
    y += Inches(0.35)

    label_bar(s, ML, y, CW, '6-week project structure', fill=ORANGE)
    y += Inches(0.27)

    wk_rows = [
        ['Week 1', 'Explore',
         'Onboard to BrightPath; shadow staff through the manual process at least once; interview 2–3 '
         'team members about pain points; document each step of the current workflow (who does it, how '
         'long it takes, tools used); assign team roles',
         'Current workflow documentation (step-by-step); staff interview notes; roles assigned'],
        ['Week 2', 'Explore + Build',
         'Analyze the inquiry dataset: calculate average response time, identify peak inquiry days, '
         'find bottlenecks; create a process map of the current workflow (flowchart); research '
         'automation tools (Zapier, Make, Apps Script); explore free tiers and capabilities',
         'Dataset analysis summary; process map flowchart; automation tool comparison chart'],
        ['Week 3', 'Build',
         'Design the automated workflow: map out which steps can be automated and how; choose '
         'automation tool; build the first automated trigger (form submission → spreadsheet row); '
         'test with sample data; document build process with screenshots',
         'Automation design diagram; Step 1 automation built and tested; build documentation'],
        ['Week 4', 'Build + Apply',
         'Build automation for steps 2–3 (welcome email + task creation); test full flow end-to-end; '
         'QA Lead documents any errors or edge cases; fix bugs; mid-program check-in; begin measuring '
         'time savings vs. manual process',
         'Automation steps 1–3 complete; QA test log; time savings measurement (draft); mid-program self-assessment'],
        ['Week 5', 'Apply',
         'Complete full automation prototype (all 5 steps or maximum feasible); conduct final testing '
         'with full sample dataset; measure total time savings; build before/after comparison; create '
         'presentation deck; practice presenting; incorporate supervisor feedback',
         'Complete automation prototype; before/after time comparison; final presentation deck; practice run completed'],
        ['Week 6', 'Launch + Reflect',
         'Final presentation to BrightPath leadership (15–20 min); live demo of automation prototype; '
         'present time savings data; answer leadership Q&A; receive formal evaluation; complete DYCD '
         'exit survey; team celebration; write personal reflection',
         'Final presentation delivered; live demo completed; evaluation form signed; exit survey; personal reflection'],
    ]
    col_w = [Inches(0.62), Inches(0.85), Inches(3.8), Inches(2.13)]
    make_table(s, ML, y, CW, Inches(3.9), ['Week', 'Phase', 'Activities', 'Deliverable'],
               wk_rows, col_widths=col_w, row_size=7.5)
    y += Inches(3.98)

    txt(s, ML, y, C1W, Inches(0.24),
        'Expected final deliverables', size=9, color=NAVY, bold=True)
    txt(s, C2X, y, C2W, Inches(0.24),
        'Key learning outcomes', size=9, color=NAVY, bold=True)
    y += Inches(0.26)

    deliverables = [
        '\u2705 Current State Process Map — detailed flowchart of the 5-step manual workflow',
        '\u2705 Automation Prototype — working automation built in Zapier/Make/Apps Script',
        '\u2705 QA Test Log — documented test cases, results, and bug fixes',
        '\u2705 Before/After Comparison — time saved, error reduction, efficiency metrics',
        '\u2705 Implementation Recommendation — steps for BrightPath to adopt the solution',
        '\u2705 Final Presentation Deck — 10–15 slides with live demo',
    ]
    for d in deliverables:
        txt(s, ML, y, C1W, Inches(0.22), d, size=7.5, color=DARK, wrap=True)
        y += Inches(0.23)

    outcomes = ['Process mapping', 'Zapier / Make', 'Google Apps Script',
                'QA testing', 'Efficiency analysis', 'Problem solving',
                'No-code tools', 'Technical documentation']
    oy = y - len(deliverables) * Inches(0.23)
    for o in outcomes:
        rect(s, C2X, oy, C2W, Inches(0.22), STRIPE if outcomes.index(o) % 2 else WHITE)
        txt(s, C2X + Pt(6), oy + Pt(2), C2W - Pt(8), Inches(0.2),
            f'\u2022  {o}', size=8, color=DARK)
        oy += Inches(0.24)

    # Prototype scope note
    sy = max(y, oy) + Inches(0.06)
    rect(s, ML, sy, CW, Inches(0.32), STRIPE)
    txt(s, ML + Pt(6), sy + Pt(4), CW - Pt(8), Inches(0.28),
        '\u26a1 Prototype Scope Note: Interns are not expected to fully deploy the automation in '
        'production. A working proof-of-concept using sample/test accounts is sufficient and '
        'demonstrates the value of the solution.',
        size=7.5, color=NAVY, wrap=True)
    sy += Inches(0.38)

    # Program contacts — navy box
    rect(s, ML, sy, CW, Inches(0.24), NAVY)
    txt(s, ML + Pt(6), sy + Pt(3), CW - Pt(8), Inches(0.22),
        'Program contacts', size=9, color=WHITE, bold=True)
    sy += Inches(0.26)
    rect(s, ML, sy, CW, Inches(0.52), LGRAY)
    txt(s, ML + Pt(6), sy + Pt(4), CW - Pt(8), Inches(0.48),
        'For DYCD SYEP support, compliance questions, or to report issues:\n'
        '\u2022  SYEP Provider — contact information available via worksite portal\n'
        '\u2022  DYCD — nyc.gov/dycd  |  1-800-246-4646',
        size=8, color=DARK, wrap=True)

    print('  Slide 14 done')

# ── Slide 15: Scenario D Intro ────────────────────────────────────────────────

def slide15(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'NYC DYCD — Scenario D: Workplace Problem Solving & App Concept Design', 15,
            'NYC DYCD — Scenario D: Workplace Problem Solving & App Concept Design')

    y = BODY_Y
    rect(s, 0, y, W, Inches(0.3), ORANGE)
    txt(s, ML, y + Pt(3), CW, Inches(0.26),
        'PROJECT SCENARIO D  —  TRACK: WORKPLACE PROBLEM SOLVING', size=9, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER)
    y += Inches(0.34)

    txt(s, ML, y, CW, Inches(0.28),
        'Workplace problem solving & app concept design', size=12, color=NAVY, bold=True)
    y += Inches(0.28)
    txt(s, ML, y, CW, Inches(0.2),
        'Host Company: Any Industry (Retail, Nonprofit, Operations, Services)  |  Duration: 6 Weeks  |  Team Size: 1–5 Interns',
        size=8, color=GRAY)
    y += Inches(0.26)

    lx = ML
    ly = y

    label_bar(s, lx, ly, C1W, 'Company background')
    ly += Inches(0.26)
    txt(s, lx, ly, C1W, Inches(0.65),
        'Your organization operates in a fast-paced environment where staff face real challenges, '
        'from customer service issues to operational inefficiencies, that often go unaddressed. '
        'Every workplace has problems better systems could address.',
        size=8, color=DARK, wrap=True)
    ly += Inches(0.7)

    txt(s, lx, ly, C1W, Inches(0.6),
        'Your intern team will:\n'
        '\u2022  Identify a real workplace challenge\n'
        '\u2022  Conduct structured research\n'
        '\u2022  Design a solution (app/tool/system concept)\n'
        '\u2022  Present a clear proposal',
        size=8, color=DARK, wrap=True)
    ly += Inches(0.68)

    rect(s, lx, ly, C1W, Inches(0.45), STRIPE)
    txt(s, lx + Pt(6), ly + Pt(4), C1W - Pt(8), Inches(0.4),
        '\u2605 Intern Mission: Research a real workplace challenge, gather input from staff or '
        'customers, and design a technology-based solution that improves how work gets done.',
        size=8, color=NAVY, wrap=True)
    ly += Inches(0.52)

    label_bar(s, lx, ly, C1W, 'Research approach')
    ly += Inches(0.26)
    research = [
        'Staff interviews (2–5 employees)',
        'Customer observations or feedback (if appropriate)',
        'Simple surveys (paper or Google Forms)',
        'Observation of daily workflows',
    ]
    for r in research:
        txt(s, lx + Pt(8), ly, C1W - Pt(10), Inches(0.2),
            f'\u2022  {r}', size=8, color=DARK)
        ly += Inches(0.21)

    ly += Inches(0.08)
    label_bar(s, lx, ly, C1W, 'Example problem areas')
    ly += Inches(0.26)
    problems = [
        'Long customer wait times (retail, service)',
        'Inefficient scheduling or communication',
        'Inventory tracking challenges',
        'Customer feedback not being captured',
        'Manual processes that could be streamlined',
    ]
    for p in problems:
        txt(s, lx + Pt(8), ly, C1W - Pt(10), Inches(0.2),
            f'\u2022  {p}', size=8, color=DARK)
        ly += Inches(0.21)

    # Example outcome
    rect(s, lx, ly + Inches(0.06), C1W, Inches(0.45), STRIPE)
    txt(s, lx + Pt(6), ly + Inches(0.1), C1W - Pt(8), Inches(0.4),
        'Example Outcome: Retail store interns found employees spent 15–20 min/shift tracking '
        'inventory manually. Their app concept cut this to under 5 minutes, improving accuracy '
        'across teams.',
        size=7.5, color=DARK, italic=True, wrap=True)

    # Right column
    rx = C2X
    ry = y

    label_bar(s, rx, ry, C2W, 'Team roles')
    ry += Inches(0.26)

    # Exact 5 roles, Solution Designer = confirmed wording
    roles_rows = [
        ['Project\nManager',   'Coordinates timeline, assigns tasks, and tracks progress'],
        ['Research\nLead',     'Conducts interviews/surveys; documents findings'],
        ['Process\nAnalyst',   'Maps current workflow; identifies inefficiencies'],
        ['Solution\nDesigner', 'Designs app/tool concept; sketches user flow'],
        ['Presentation\nLead', 'Builds final deck and delivers pitch'],
    ]
    rcw = [Inches(1.15), Inches(2.57)]
    make_table(s, rx, ry, C2W, Inches(1.95), ['Role', 'Responsibilities'], roles_rows,
               col_widths=rcw, row_size=7.5)
    ry += Inches(2.02)

    label_bar(s, rx, ry, C2W, 'Solution output (no coding required)')
    ry += Inches(0.26)
    outputs = [
        'Mobile app concept',
        'Digital workflow or tracking tool',
        'Communication or scheduling system',
        'Operations dashboard',
    ]
    for o in outputs:
        txt(s, rx + Pt(8), ry, C2W - Pt(10), Inches(0.2),
            f'\u2022  {o}', size=8, color=DARK)
        ry += Inches(0.21)

    ry += Inches(0.08)
    label_bar(s, rx, ry, C2W, 'Tools (no coding required)')
    ry += Inches(0.26)
    tools = ['Paper sketches', 'Google Slides', 'Canva', 'Figma (optional)']
    for t in tools:
        txt(s, rx + Pt(8), ry, C2W - Pt(10), Inches(0.2),
            f'\u2022  {t}', size=8, color=DARK)
        ry += Inches(0.21)

    ry += Inches(0.08)
    rect(s, rx, ry, C2W, Inches(0.42), STRIPE)
    txt(s, rx + Pt(6), ry + Pt(4), C2W - Pt(8), Inches(0.38),
        '\u2605 No Coding Required: Interns are not expected to build a working app. '
        'A well-researched concept, wireframe, or prototype is sufficient to demonstrate '
        'the value of the solution.',
        size=8, color=NAVY, wrap=True)

    print('  Slide 15 done')

# ── Slide 16: Scenario D — 6-Week Structure ───────────────────────────────────

def slide16(prs):
    s = new_slide(prs)
    hdr_ftr(s, 'NYC DYCD — Scenario D: Workplace Problem Solving & App Concept Design', 16,
            'NYC DYCD — Scenario D: Workplace Problem Solving & App Concept Design')

    y = BODY_Y
    txt(s, ML, y, CW, Inches(0.32),
        'Scenario D: Workplace problem solving — 6-week structure & deliverables',
        size=12, color=NAVY, bold=True)
    y += Inches(0.35)

    # Title "6-week project structure" in ORANGE (confirmed fix)
    label_bar(s, ML, y, CW, '6-week project structure', fill=ORANGE)
    y += Inches(0.27)

    wk_rows = [
        ['Week 1', 'Explore',
         'Observe company operations; identify potential challenges; assign team roles; set up project plan',
         'List of 2–3 workplace challenges'],
        ['Week 2', 'Explore + Build',
         'Interview staff; conduct surveys; analyze findings; draft problem statement; select focus area',
         'Problem statement + research summary'],
        ['Week 3', 'Build',
         'Map current workflow; identify inefficiencies; brainstorm solution ideas; select concept with supervisor',
         'Workflow map + solution concept draft'],
        ['Week 4', 'Build + Apply',
         'Design solution concept; define key features; gather staff feedback; mid-program check-in',
         'Draft solution design + feedback notes; mid-program self-assessment'],
        ['Week 5', 'Apply',
         'Refine solution; build before/after comparison; create presentation deck; practice with supervisor',
         'Final solution concept + presentation draft'],
        ['Week 6', 'Launch + Reflect',
         'Present to employer (10–15 min); receive feedback; complete exit survey; team debrief; reflect on learning',
         'Final presentation + personal reflection'],
    ]
    col_w = [Inches(0.62), Inches(0.85), Inches(3.8), Inches(2.13)]
    make_table(s, ML, y, CW, Inches(3.5), ['Week', 'Phase', 'Activities', 'Deliverable'],
               wk_rows, col_widths=col_w, row_size=8)
    y += Inches(3.58)

    txt(s, ML, y, C1W, Inches(0.24),
        'Expected final deliverables', size=9, color=NAVY, bold=True)
    txt(s, C2X, y, C2W, Inches(0.24),
        'Key learning outcomes', size=9, color=NAVY, bold=True)
    y += Inches(0.26)

    deliverables = [
        '\u2022  Problem Statement — clearly defined challenge with supporting evidence',
        '\u2022  Research Summary — interviews, observations, insights',
        '\u2022  Workflow Map — current process visualization',
        '\u2022  Solution Concept — app / tool / system design',
        '\u2022  Before/After Comparison — time saved or impact shown',
        '\u2022  Final Presentation — 10–15 slides',
    ]
    for d in deliverables:
        txt(s, ML, y, C1W, Inches(0.22), d, size=7.5, color=DARK, wrap=True)
        y += Inches(0.23)

    outcomes = ['Problem solving', 'User research', 'Process analysis',
                'Critical thinking', 'Communication', 'Design thinking', 'Workplace collaboration']
    oy = y - len(deliverables) * Inches(0.23)
    for o in outcomes:
        rect(s, C2X, oy, C2W, Inches(0.22), STRIPE if outcomes.index(o) % 2 else WHITE)
        txt(s, C2X + Pt(6), oy + Pt(2), C2W - Pt(8), Inches(0.2),
            f'\u2022  {o}', size=8, color=DARK)
        oy += Inches(0.24)

    # Prototype scope note
    sy = max(y, oy) + Inches(0.06)
    rect(s, ML, sy, CW, Inches(0.32), STRIPE)
    txt(s, ML + Pt(6), sy + Pt(4), CW - Pt(8), Inches(0.28),
        'Prototype Scope Note: Interns are not expected to build a working app or deploy a live system. '
        'A researched concept or prototype is sufficient.',
        size=7.5, color=NAVY, wrap=True)
    sy += Inches(0.38)

    # Program contacts — white background style (Slide 16 variant)
    rect(s, ML, sy, CW, Inches(0.24), NAVY)
    txt(s, ML + Pt(6), sy + Pt(3), CW - Pt(8), Inches(0.22),
        'Program contacts', size=9, color=WHITE, bold=True)
    sy += Inches(0.26)
    rect(s, ML, sy, CW, Inches(0.62), WHITE)
    txt(s, ML + Pt(6), sy + Pt(4), CW - Pt(8), Inches(0.58),
        'For DYCD SYEP support, compliance questions, or to report issues:\n'
        '\u2022  SYEP Provider — contact information available via worksite portal\n'
        '\u2022  DYCD — nyc.gov/dycd  |  1-800-246-4646\n'
        '\u2022  Employer Feedback Form — ______________________________',
        size=8, color=DARK, wrap=True)

    print('  Slide 16 done')

# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print('Building slides...')
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    slide7(prs)
    slide8(prs)
    slide9(prs)
    slide10(prs)
    slide11(prs)
    slide12(prs)
    slide13(prs)
    slide14(prs)
    slide15(prs)
    slide16(prs)

    prs.save(OUT)
    print(f'\nSaved: {OUT}')
    print(f'Total slides: {len(prs.slides)}')

if __name__ == '__main__':
    main()
